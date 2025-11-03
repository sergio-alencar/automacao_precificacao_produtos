# src/generators/pptx_generator.py

import logging
import os
from typing import Optional, List, Tuple
from pptx import Presentation
from pptx.util import Inches

from . import pptx_helpers
from . import pptx_table_builder


logger = logging.getLogger(__name__)


def generate_pptx_presentation(
    template_path: str,
    municipal_name: str,
    state_uf: str,
    product_results: list[tuple[str, Optional[float]]],
    output_filename: str = "Apresentacao_MSL_Preenchida.pptx",
    export_to_pdf: bool = True,
):
    if not os.path.exists(template_path):
        logger.error(f"PPTX template not found: {template_path}")
        raise FileNotFoundError(f"PPTX template not found at {template_path}")

    try:
        prs = Presentation(template_path)
        logger.info(f"Template '{template_path}' loaded.")
    except Exception as e:
        logger.error(f"Error loading PPTX template: {e}")
        raise

    slide_2 = prs.slides[1]

    total_sum = sum(val for _, val in product_results if val is not None)

    total_text_placeholder = "{{TOTAL_FINAL}}"
    formatted_total_text = f"R$ {pptx_helpers.format_total_summary(total_sum)}"

    pptx_helpers.find_and_replace_text(slide_2, total_text_placeholder, formatted_total_text)
    logger.info(f"Total summary text replaced on slide 2: '{formatted_total_text}'.")

    pptx_helpers.remove_all_tables_from_slide(slide_2)

    table_left = Inches(1.25)
    table_top = Inches(2.0)
    table_width = Inches(8.5)
    table_height = Inches(5.0)

    pptx_table_builder.create_results_table(
        slide_2,
        product_results,
        municipal_name,
        state_uf,
        table_left,
        table_top,
        table_width,
        table_height,
    )

    slide_3 = prs.slides[2]

    old_municipal_uf_text = "{{MUNICIPIO_UF}}"
    new_municipal_uf_text = f"{municipal_name} - {state_uf}"

    pptx_helpers.find_and_replace_text(slide_3, old_municipal_uf_text, new_municipal_uf_text)
    logger.info(f"Municipality/UF replaced on slide 3: '{new_municipal_uf_text}'.")

    try:
        output_pptx_path = output_filename
        prs.save(output_pptx_path)
        logger.info(f"Filled PPTX presentation saved to: '{output_pptx_path}'.")
    except Exception as e:
        logger.error(f"Error saving PPTX presentation: {e}")
        raise

    if export_to_pdf:
        logger.warning(
            "Direct PPTX to PDF conversion is not supported by `python-pptx`.\nPPTX file was generated. For PDF, convert manually or use an external API/tool."
        )
