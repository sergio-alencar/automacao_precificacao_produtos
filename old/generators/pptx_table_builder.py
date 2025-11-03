# src/generators/pptx_table_builder.py

import logging
from typing import Optional, List, Tuple
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from .pptx_helpers import format_currency


logger = logging.getLogger(__name__)


def create_results_table(
    slide: Slide,
    product_results: List[Tuple[str, Optional[float]]],
    municipal_name: str,
    state_uf: str,
    x: Inches,
    y: Inches,
    width: Inches,
    height: Inches,
):
    rows = len(product_results) + 2
    cols = 2

    table_shape = slide.shapes.add_table(
        rows,
        cols,
        x,
        y,
        width,
        height,
    ).table

    table_shape.cell(0, 0).text = "Produto"
    table_shape.cell(0, 1).text = "Estimativa"

    total_sum = 0.0
    for i, (product_name, value) in enumerate(product_results):
        row_idx = i + 1
        table_shape.cell(row_idx, 0).text = product_name
        if value is not None:
            table_shape.cell(row_idx, 1).text = format_currency(value)
            total_sum += value
        else:
            table_shape.cell(row_idx, 1).text = "Não aplicável"

    last_row_idx = rows - 1
    table_shape.cell(last_row_idx, 0).text = "TOTAL GERAL"
    table_shape.cell(last_row_idx, 1).text = format_currency(total_sum)

    for r in range(rows):
        for c in range(cols):
            cell = table_shape.cell(r, c)
            text_frame = cell.text_frame
            text_frame.margin_left = Pt(5)
            text_frame.margin_right = Pt(5)
            text_frame.margin_top = Pt(5)
            text_frame.margin_bottom = Pt(5)

            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.name = "Arial"
            p.font.color.rgb = RGBColor(0, 0, 0)

            if c == 1:
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT

            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            if r == 0:
                p.font.bold = True
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(100, 100, 100)
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

        if r == last_row_idx:
            for c_total in range(cols):
                cell = table_shape.cell(r, c_total)
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0, 0, 0)
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    table_shape.columns[0].width = Inches(4.5)
    table_shape.columns[1].width = Inches(2.5)

    left = x
    top = y - Pt(30)
    width_title = width
    height_title = Pt(20)

    title_shape = slide.shapes.add_textbox(left, top, width_title, height_title)  # type: ignore
    text_frame_title = title_shape.text_frame
    text_frame_title.text = f"Estimativas de Potenciais de Recuperação - {municipal_name}/{state_uf}"
    text_frame_title.paragraphs[0].font.name = "Arial"
    text_frame_title.paragraphs[0].font.size = Pt(12)
    text_frame_title.paragraphs[0].font.bold = True
    text_frame_title.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    logger.info("Results table created on slide.")
