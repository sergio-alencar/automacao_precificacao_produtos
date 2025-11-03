# src/generators/pptx_helpers.py

import locale
import logging
from pptx.slide import Slide

logger = logging.getLogger(__name__)

try:
    locale.setlocale(locale.LC_ALL, "pt_BR.UTF-8")
except locale.Error:
    try:
        locale.setlocale(locale.LC_ALL, "Portuguese_Brazil.1252")
    except locale.Error:
        logger.warning("pt_BR locale not found. Using system standard locale.")
        locale.setlocale(locale.LC_ALL, "")


def format_currency(value: float) -> str:
    return locale.format_string("R$ %.2f", value, grouping=True)


def format_total_summary(value: float) -> str:
    if value < 1_000_000:
        return f"{value/1_000:.0f} MIL"
    else:
        return f"{value/1_000_000:.0f} MILHÕES"


def find_and_replace_text(slide: Slide, old_text: str, new_text: str):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text_frame = shape.text_frame  # type: ignore
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
                    logger.debug(f"Replaced '{old_text}' with '{new_text}' in shape: {shape.name}")


def remove_all_tables_from_slide(slide: Slide):
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.has_table:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
        logger.debug("Table removed from slide.")
